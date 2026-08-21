"""生成 GOAI 2026 初赛 PPT 第一版。

内容与 hackathon/goai-2026-submission/PPT-内容稿.md 对齐。
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SLIDES = [
    {
        "title": "WeFinance 智能记账与经营流水助手",
        "bullets": [
            "赛道：无界应用（Boundless Agents）｜AI+金融",
            "一句话定位：拍照上传账单，GPT-4o Vision 一步完成结构化提取，结合异常检测与可解释理财分析，为个体户和自由职业者提供经营流水画像",
        ],
    },
    {
        "title": "问题与场景",
        "bullets": [
            "三类用户：个体户、自由职业者、习惯手工记账的普通用户",
            "痛点：传统 OCR 对票据识别率低（PaddleOCR 在合成账单上曾实测 0%）；异常消费难及时发现；理财建议千篇一律、不基于真实数据",
            "场景：拍照上传账单 -> 自动结构化 -> 消费洞察 -> AI 顾问问答 -> 可解释投资建议",
        ],
    },
    {
        "title": "方案总览",
        "bullets": [
            "核心闭环：账单上传 -> Vision LLM 结构化提取 -> 异常检测 -> 消费洞察 -> 对话式顾问 -> 可解释理财报告",
            "一步式 OCR：GPT-4o Vision 直接输出结构化交易数据，替代 OCR+文本理解两步管线",
            "新增经营流水画像视图：把已有交易数据聚合为交易对手集中度、收支趋势、异常预警",
        ],
    },
    {
        "title": "Agent 能力",
        "bullets": [
            "任务输入：账单图片、用户预算、聊天问题",
            "意图理解与任务规划：LangChain AgentExecutor，识别查预算/查分类/查具体类目等意图",
            "能力调用：query_budget、query_spending、query_category 三类工具，结合对话记忆多轮交互",
            "结果交付：结构化交易记录、异常预警卡片、经营画像聚合视图（交易对手集中度/收支趋势）、可解释推理链（理财建议附带 rationale_steps）",
            "验证与反馈：异常检测支持用户确认/标记误报反馈闭环",
        ],
    },
    {
        "title": "技术与工程",
        "bullets": [
            "Vision OCR：GPT-4o Vision，温度 0.0，30 秒超时保护",
            "异常检测：z-score 统计方法，按样本量自适应调整灵敏度和阈值",
            "可解释推荐：输出带因果推理步骤的理财建议",
            "多数服务层通过 @safe_call 装饰器提供超时保护和优雅降级；Chat 走独立的重试+退避机制",
            "技术栈：Streamlit + LangChain + OpenAI 兼容 API + Pydantic + Plotly",
        ],
    },
    {
        "title": "Demo 演示",
        "bullets": [
            "本地实测（2026-08-08，通过）：种子数据走完消费洞察 -> 经营流水画像新视图，数字随数据变化，非写死样例",
            "线上实测（2026-08-11，通过）：真实上传样例账单触发线上 OCR 调用，4/4 笔交易全部正确识别（商户/日期/金额/分类）",
            "过程中定位到服务商侧模型下线（Qwen2-VL-72B-Instruct 已 deprecated）导致的持续 403，切换至当前活跃模型后修复，非代码缺陷",
        ],
    },
    {
        "title": "数据与合规",
        "bullets": [
            "账单图片仅内存中 base64 编码后调用 Vision API，不落盘保存原图",
            "交易数据存于 session state，无数据库；如需跨会话保留，本地 JSON 持久化（默认不加密）",
            "商业 API 依赖如实披露：OpenAI 兼容接口，密钥通过 Streamlit Secrets 管理，不提交仓库",
        ],
    },
    {
        "title": "AI+金融边界声明",
        "bullets": [
            "理财建议明确声明分析与教育性内容，不构成投资/信贷/保险决策依据，不替代持牌机构判断",
            "不点名具体基金、券商或投资平台，不将执行参数（金额/时点/价位）当作确定性指令下达",
            "异常检测、消费分析基于用户真实数据，不做身份核验，不采集身份证号/银行卡号等敏感字段",
        ],
    },
    {
        "title": "开放与复用",
        "bullets": [
            "MIT 协议开源，中英双语 README 和文档",
            "票据识别、推荐引擎为独立服务层，可拆分复用于其他记账类应用",
            "GitHub 128 星、16 fork（2026-08-08 实测）",
        ],
    },
    {
        "title": "迭代与落地",
        "bullets": [
            "近期：线上部署问题（Python 版本、OCR 模型下线）已修复并复测通过；补齐 tests/ 测试套件",
            "中期：经营流水画像扩展更多维度（如按周期对比、导出报表）",
            "长期：探索个体户/自由职业者报税辅助场景",
        ],
    },
    {
        "title": "原有基础与本次新增贡献",
        "bullets": [
            "原有基础：2025 年 11 月开发的完整产品，票据识别、消费分析、AI 顾问、投资建议四大核心链路已上线运行，GitHub 128 星",
            "本次新增：经营流水画像视图（交易对手集中度、收支趋势聚合）；修复 generate_detailed_report 的金融合规边界",
            "实测中顺带发现并修复四个真实工程缺陷（st.dataframe 参数兼容崩溃、pyarrow/numpy ABI 不兼容导致的原生崩溃、线上 Python 版本配置错误、财务健康卡片 HTML 渲染为纯文本）；GOAI 参赛材料",
        ],
    },
    {
        "title": "结尾",
        "bullets": [
            "主张：WeFinance 让拍照记账直接变成经营流水画像，服务被传统财务软件忽视的个体户和自由职业者",
            "代码与文档全部开源，可现场演示",
            "联系方式与项目链接",
        ],
    },
]


def build(output_path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    brand_dark = RGBColor(0x0F, 0x17, 0x2A)
    brand_primary = RGBColor(0x0D, 0x94, 0x88)
    brand_accent = RGBColor(0xD4, 0xAF, 0x37)
    text_dark = RGBColor(0x1E, 0x29, 0x3B)

    def add_background(slide, color):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

    def add_brand_bar(slide):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            prs.slide_width,
            Inches(0.12),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = brand_primary
        bar.line.fill.background()

    # 封面
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(cover, brand_dark)
    title_box = cover.shapes.add_textbox(
        Inches(1), Inches(2.4), Inches(11.3), Inches(1.4)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_run = title_p.add_run()
    title_run.text = SLIDES[0]["title"]
    title_run.font.size = Pt(40)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0xF8, 0xFA, 0xFC)

    subtitle_box = cover.shapes.add_textbox(
        Inches(1), Inches(4.1), Inches(11.3), Inches(1.6)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    for i, line in enumerate(SLIDES[0]["bullets"]):
        p = subtitle_frame.paragraphs[0] if i == 0 else subtitle_frame.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.size = Pt(20)
        run.font.color.rgb = brand_accent

    for slide_data in SLIDES[1:]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide, RGBColor(0xFF, 0xFF, 0xFF))
        add_brand_bar(slide)

        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.35), Inches(12), Inches(0.9)
        )
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_run = title_p.add_run()
        title_run.text = slide_data["title"]
        title_run.font.size = Pt(32)
        title_run.font.bold = True
        title_run.font.color.rgb = brand_primary

        body_box = slide.shapes.add_textbox(
            Inches(0.9), Inches(1.45), Inches(11.6), Inches(3.6)
        )
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        for i, bullet in enumerate(slide_data["bullets"]):
            p = body_frame.paragraphs[0] if i == 0 else body_frame.add_paragraph()
            p.space_after = Pt(14)
            run = p.add_run()
            run.text = bullet
            run.font.size = Pt(18)
            run.font.color.rgb = text_dark

        if slide_data["title"] == "Demo 演示":
            screenshot = Path(
                "hackathon/goai-2026-submission/assets/demo-business-profile-detail.jpg"
            )
            if screenshot.exists():
                slide.shapes.add_picture(
                    str(screenshot), Inches(1.0), Inches(5.15), width=Inches(6.4)
                )

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"已生成: {output_path}")


if __name__ == "__main__":
    build("hackathon/goai-2026-submission/WeFinance-GOAI2026-初赛PPT-v1.pptx")
